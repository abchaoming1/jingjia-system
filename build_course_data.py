from __future__ import annotations

import json
import re
import xml.etree.ElementTree as ET
from collections import Counter
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation


ROOT = Path(__file__).resolve().parent
BASICS_PPT = Path(r"C:\Users\abcha\Desktop\集合竞价交易认知与实战要点.pptx")
CASES_PPT = Path(r"C:\Users\abcha\Desktop\接力战法 - 集合竞价案例教学.pptx")
OUTPUT_JS = ROOT / "assets" / "course-data.js"

P_NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
P14_NS = {"p14": "http://schemas.microsoft.com/office/powerpoint/2010/main"}


def asset_version(path: Path) -> str:
    return str(int(path.stat().st_mtime))


def clean_text(value: str) -> str:
    value = (
        value.replace("\xa0", " ")
        .replace("\u3000", " ")
        .replace("•", "-")
        .replace("●", "-")
        .replace("▪", "-")
        .replace("\r", "\n")
    )

    lines: list[str] = []
    for raw_line in value.splitlines():
        line = re.sub(r"\s+", " ", raw_line).strip()
        if not line:
            continue
        if line[0] in {"-", "•", "●", "▪", "·"}:
            line = "- " + line[1:].strip()
        lines.append(line)
    return " / ".join(lines)


def extract_text_blocks(slide) -> list[str]:
    seen: set[str] = set()
    blocks: list[str] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = clean_text(shape.text)
        if not text or is_progress_marker(text) or text in seen:
            continue
        seen.add(text)
        blocks.append(text)
    return blocks


def is_progress_marker(value: str) -> bool:
    return bool(re.fullmatch(r"\d+\s*/\s*\d+", value))


def split_points(text: str) -> list[str]:
    if not text:
        return []

    normalized = (
        text.replace(" / - ", "\n- ")
        .replace(" /-", "\n- ")
        .replace("/-", "\n- ")
        .replace(" / ", "\n")
    )

    points: list[str] = []
    for raw_line in normalized.splitlines():
        line = raw_line.strip(" /")
        if not line:
            continue
        if line.startswith("-"):
            line = line[1:].strip()
        if line and line not in points:
            points.append(line)
    return points


def slugify(value: str) -> str:
    slug = re.sub(r"[^\w\u4e00-\u9fff]+", "-", value.strip().lower(), flags=re.UNICODE)
    slug = slug.strip("-")
    return slug or "section"


def find_first_cjk(blocks: list[str]) -> str:
    for block in blocks:
        if re.search(r"[\u4e00-\u9fff]", block):
            return block
    return ""


def non_trivial_blocks(blocks: list[str]) -> list[str]:
    ignored = {"COURSE CHAPTER", "STRUCTURE", "COURSE HANDOUT"}
    return [
        block
        for block in blocks
        if block not in ignored and not re.fullmatch(r"CHAPTER\s+\d+", block) and not re.fullmatch(r"\d+", block)
    ]


def extract_basics_slides(ppt_path: Path) -> list[dict]:
    presentation = Presentation(str(ppt_path))
    slides: list[dict] = []

    for number, slide in enumerate(presentation.slides, start=1):
        blocks = extract_text_blocks(slide)
        eyebrow = blocks[0] if blocks else f"SLIDE {number:02d}"
        tail = non_trivial_blocks(blocks[1:])
        title = find_first_cjk(tail) or (tail[0] if tail else eyebrow)
        point_limit = 4

        subtitle = ""
        for block in tail:
            if block != title:
                subtitle = block
                break

        points: list[str] = []
        if eyebrow.startswith("REVIEW TEMPLATE"):
            points = tail[2:]
            point_limit = len(points)

        for block in tail:
            if points:
                break
            if block in {title, subtitle}:
                continue
            parsed = split_points(block)
            if len(parsed) >= 2:
                points = parsed[:5]
                break

        if not points:
            points = [block for block in tail if block not in {title, subtitle}][:4]

        summary = points[0] if points else subtitle or title
        slides.append(
            {
                "number": number,
                "eyebrow": eyebrow,
                "title": title,
                "subtitle": subtitle,
                "summary": summary,
                "points": points,
                "pointLimit": point_limit,
                "image": f"./assets/basics-slides/slide-{number:03d}.png?v={asset_version(BASICS_PPT)}",
                "_blocks": blocks,
            }
        )

    return slides


def build_basics_chapters(slides: list[dict]) -> list[dict]:
    chapter_markers = [index for index, slide in enumerate(slides) if slide["eyebrow"].startswith("CHAPTER")]
    chapters: list[dict] = []

    if chapter_markers and chapter_markers[0] > 0:
        intro_slides = slides[: chapter_markers[0]]
        chapters.append(
            {
                "id": "basics-00-导读与结构",
                "index": 0,
                "name": "导读与结构",
                "slideCount": len(intro_slides),
                "summary": f"本章共 {len(intro_slides)} 页，先交代课程定位、复盘目录与整套基础认知的阅读方式。",
                "slides": strip_internal_blocks(intro_slides),
            }
        )

    for order, marker_index in enumerate(chapter_markers, start=1):
        next_marker = chapter_markers[order] if order < len(chapter_markers) else len(slides)
        chunk = slides[marker_index:next_marker]
        marker_slide = chunk[0]

        chapter_name = find_first_cjk(marker_slide["_blocks"]) or marker_slide["title"]
        chapter_slides = strip_internal_blocks(chunk)
        topics = [slide["title"] for slide in chapter_slides[1:4] if slide["title"]][:3]
        summary = f"本章共 {len(chapter_slides)} 页，围绕“{chapter_name}”展开。"
        if topics:
            summary += f" 重点包括：{'、'.join(topics)}。"

        chapters.append(
            {
                "id": f"basics-{order:02d}-{slugify(chapter_name)}",
                "index": order,
                "name": chapter_name,
                "slideCount": len(chapter_slides),
                "summary": summary,
                "slides": chapter_slides,
            }
        )

    return chapters


def strip_internal_blocks(slides: list[dict]) -> list[dict]:
    cleaned: list[dict] = []
    for slide in slides:
        item = dict(slide)
        item.pop("_blocks", None)
        cleaned.append(item)
    return cleaned


def build_basics_payload() -> dict:
    slides = extract_basics_slides(BASICS_PPT)
    chapters = build_basics_chapters(slides)

    return {
        "title": "集合竞价交易认知与实战要点",
        "sourceFile": str(BASICS_PPT),
        "chapterCount": len(chapters),
        "slideCount": len(slides),
        "chapters": chapters,
    }


def parse_market_context(texts: list[str]) -> tuple[str, str]:
    for text in reversed(texts):
        if any(token in text for token in ("题材", "龙头", "中位", "补涨", "中军", "老龙")):
            if "：" in text:
                market_context, stock_name = text.rsplit("：", 1)
                return market_context.strip(), stock_name.strip()
            if ":" in text:
                market_context, stock_name = text.rsplit(":", 1)
                return market_context.strip(), stock_name.strip()
            return text, ""
    return "", ""


def infer_case_tags(title: str, section_name: str, market_context: str) -> list[str]:
    candidates = [
        ("一字板", "一字"),
        ("弱转强", "弱转强"),
        ("抢筹", "抢筹"),
        ("爆量", "爆量"),
        ("长阳线", "长阳线"),
        ("锤头", "锤头"),
        ("大分歧", "大分歧"),
        ("小分歧", "小分歧"),
        ("承接", "承接"),
        ("中军", "中军"),
    ]
    source = " | ".join([title, section_name, market_context])
    tags = [label for keyword, label in candidates if keyword in source]
    return tags[:4]


def parse_theme_label(market_context: str) -> str | None:
    if not market_context:
        return None
    match = re.search(r"（([^）]+)）", market_context)
    if match:
        return match.group(1).strip()
    if "-" in market_context:
        return market_context.split("-", 1)[0].strip()
    return market_context.strip()


def extract_case_slides(ppt_path: Path) -> list[dict]:
    presentation = Presentation(str(ppt_path))
    slides: list[dict] = []

    for number, slide in enumerate(presentation.slides, start=1):
        texts = extract_text_blocks(slide)
        title = texts[0] if texts else f"案例 {number}"
        extra_texts = texts[1:]
        market_context, stock_name = parse_market_context(extra_texts)

        bullet_source = ""
        notes: list[str] = []
        for block in extra_texts:
            if block == market_context or block == stock_name:
                continue
            points = split_points(block)
            if len(points) >= 2 and not bullet_source:
                bullet_source = block
            else:
                notes.append(block)

        takeaways = split_points(bullet_source)[:5]
        if not takeaways and extra_texts:
            takeaways = split_points(" / ".join(extra_texts))[:5]

        summary = takeaways[0] if takeaways else (extra_texts[0] if extra_texts else "")
        slides.append(
            {
                "number": number,
                "title": title,
                "summary": summary,
                "takeaways": takeaways,
                "marketContext": market_context,
                "stockName": stock_name,
                "notes": notes[:3],
                "image": f"./assets/slides/slide-{number:03d}.png?v={asset_version(CASES_PPT)}",
            }
        )

    return slides


def extract_case_sections(ppt_path: Path) -> list[dict]:
    with ZipFile(ppt_path) as archive:
        presentation_xml = ET.fromstring(archive.read("ppt/presentation.xml"))

    slide_id_map: dict[str, int] = {}
    for position, slide_tag in enumerate(presentation_xml.findall("./p:sldIdLst/p:sldId", P_NS), start=1):
        slide_id_map[slide_tag.attrib["id"]] = position

    sections: list[dict] = []
    for section_tag in presentation_xml.findall(".//p14:sectionLst/p14:section", P14_NS):
        slide_numbers = []
        for slide_id_tag in section_tag.findall("./p14:sldIdLst/p14:sldId", P14_NS):
            number = slide_id_map.get(slide_id_tag.attrib["id"])
            if number:
                slide_numbers.append(number)
        if slide_numbers:
            sections.append(
                {
                    "name": section_tag.attrib["name"],
                    "slideNumbers": sorted(slide_numbers),
                }
            )

    return sorted(sections, key=lambda item: item["slideNumbers"][0])


def build_case_section_summary(name: str, cases: list[dict]) -> str:
    stock_names = [case["stockName"] for case in cases if case["stockName"]][:3]

    if "转分歧" in name:
        focus = "重点看强一致之后，竞价爆量、抛压释放和次日承接是否失衡。"
    elif "继续加速" in name or "一致加速" in name:
        focus = "重点看强势加速阶段，试盘、惜售和尾段抢筹是否同步强化。"
    elif "弱转强" in name:
        focus = "重点看竞价偏弱时，盘中承接如何完成从弱到强的切换。"
    elif "烂板" in name or "易炸板" in name:
        focus = "重点看大分歧、高抛压环境下，回封质量和炸板风险的差别。"
    elif "中军" in name:
        focus = "重点看大市值中军的竞价节奏与普通连板票的执行差异。"
    else:
        focus = "重点看试盘、抛压、分歧、承接和竞昨比之间的协同关系。"

    if stock_names:
        focus += f" 代表个股：{'、'.join(stock_names)}。"

    return f"本节收录 {len(cases)} 页案例，{focus}"


def build_cases_payload() -> dict:
    slides = extract_case_slides(CASES_PPT)
    slides_by_number = {slide["number"]: slide for slide in slides}
    raw_sections = extract_case_sections(CASES_PPT)

    sections: list[dict] = []
    for index, raw_section in enumerate(raw_sections, start=1):
        cases = [slides_by_number[number] for number in raw_section["slideNumbers"] if number in slides_by_number]
        name = raw_section["name"]
        market_labels = [case["marketContext"] for case in cases if case["marketContext"]]

        sections.append(
            {
                "id": f"cases-{index:02d}-{slugify(name)}",
                "index": index,
                "name": name,
                "slideCount": len(cases),
                "summary": build_case_section_summary(name, cases),
                "keywords": infer_case_tags(name, name, " | ".join(market_labels)),
                "cases": cases,
            }
        )

    theme_counter = Counter()
    for slide in slides:
        theme = parse_theme_label(slide["marketContext"])
        if theme:
            theme_counter[theme] += 1

    return {
        "title": "接力战法 - 集合竞价案例教学",
        "sourceFile": str(CASES_PPT),
        "sectionCount": len(sections),
        "slideCount": len(slides),
        "topThemes": [{"name": name, "count": count} for name, count in theme_counter.most_common(8)],
        "sections": sections,
    }


def build_payload() -> dict:
    basics = build_basics_payload()
    cases = build_cases_payload()
    return {
        "title": "集合竞价课程双模块",
        "basics": basics,
        "cases": cases,
    }


def main() -> None:
    for path in (BASICS_PPT, CASES_PPT):
        if not path.exists():
            raise SystemExit(f"PPT not found: {path}")

    payload = build_payload()
    OUTPUT_JS.parent.mkdir(parents=True, exist_ok=True)
    OUTPUT_JS.write_text(
        "window.courseData = " + json.dumps(payload, ensure_ascii=False, indent=2) + ";\n",
        encoding="utf-8",
    )
    print(f"Generated {OUTPUT_JS}")


if __name__ == "__main__":
    main()
